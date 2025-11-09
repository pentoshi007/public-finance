#!/usr/bin/env python3
"""
Convert HTML presentation to PPTX using Playwright for screenshots
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import asyncio
from playwright.async_api import async_playwright

async def capture_slides(html_file, output_dir="slides_screenshots"):
    """Capture each slide as a screenshot"""
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={'width': 1280, 'height': 720})
        
        # Load the HTML file
        await page.goto(f'file://{os.path.abspath(html_file)}')
        await page.wait_for_load_state('networkidle')
        
        # Find all slide wrappers
        slides = await page.query_selector_all('.slide-wrapper')
        print(f"Found {len(slides)} slides")
        
        for i, slide in enumerate(slides):
            # Scroll to slide
            await slide.scroll_into_view_if_needed()
            await asyncio.sleep(0.5)  # Wait for any animations
            
            # Take screenshot
            screenshot_path = os.path.join(output_dir, f'slide_{i+1}.png')
            await slide.screenshot(path=screenshot_path)
            print(f"Captured slide {i+1}")
        
        await browser.close()
    
    return len(slides)

def create_pptx_from_images(num_slides, output_file="presentation.pptx", screenshots_dir="slides_screenshots"):
    """Create PPTX from screenshot images"""
    prs = Presentation()
    
    # Set slide size to 16:9 (standard presentation size)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    for i in range(1, num_slides + 1):
        # Add blank slide
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add image
        img_path = os.path.join(screenshots_dir, f'slide_{i}.png')
        if os.path.exists(img_path):
            # Add image to fill the entire slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            print(f"Added slide {i} to PPTX")
    
    prs.save(output_file)
    print(f"\nPresentation saved as: {output_file}")

async def main():
    html_file = "ppt.html"
    output_pptx = "presentation.pptx"
    screenshots_dir = "slides_screenshots"
    
    print("Step 1: Capturing slides as screenshots...")
    num_slides = await capture_slides(html_file, screenshots_dir)
    
    print("\nStep 2: Creating PPTX from screenshots...")
    create_pptx_from_images(num_slides, output_pptx, screenshots_dir)
    
    print("\n✅ Conversion complete!")
    print(f"📄 Output file: {output_pptx}")
    print(f"🖼️  Screenshots saved in: {screenshots_dir}/")

if __name__ == "__main__":
    asyncio.run(main())

